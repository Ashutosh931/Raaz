"""
PowerPoint Presentation Generator for Final Year Project:
AI-Based Smart Energy Consumption Prediction and Optimization System
Generates a professional 15-slide PowerPoint deck (presentation.pptx).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_filename="presentation.pptx"):
    prs = Presentation()
    # Set 16:9 widescreen dimensions (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Color Palette
    COLOR_PRIMARY = RGBColor(16, 185, 129)    # Emerald Green
    COLOR_SECONDARY = RGBColor(14, 165, 233)  # Electric Sky Blue
    COLOR_DARK = RGBColor(15, 23, 42)         # Slate 900
    COLOR_LIGHT = RGBColor(248, 250, 252)     # Off-white / light
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_GRAY = RGBColor(100, 116, 139)      # Slate 500
    COLOR_CARD_BG = RGBColor(241, 245, 249)   # Slate 100
    COLOR_ACCENT = RGBColor(245, 158, 11)     # Amber

    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="AI SMART ENERGY OPTIMIZATION SYSTEM"):
        # Header Background Accent Bar
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(11.73), Inches(0.9)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = COLOR_LIGHT
        header_shape.line.color.rgb = COLOR_PRIMARY
        header_shape.line.width = Pt(1.5)

        tf = header_shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.1)

        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_PRIMARY

        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_DARK

    def add_card(slide, left, top, width, height, title, content_items, accent_color=COLOR_PRIMARY, bg_color=COLOR_CARD_BG):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = accent_color
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)

        if title:
            p_t = tf.paragraphs[0]
            p_t.text = title
            p_t.font.size = Pt(16)
            p_t.font.bold = True
            p_t.font.color.rgb = COLOR_DARK
            p_t.space_after = Pt(10)

        for item in content_items:
            p = tf.add_paragraph()
            p.text = f"•  {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = COLOR_DARK
            p.space_after = Pt(6)

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_DARK
    bg1.line.fill.background()

    # Title Card Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(4.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True

    p_badge = tf1.paragraphs[0]
    p_badge.text = "FINAL YEAR CAPSTONE PROJECT PRESENTATION"
    p_badge.font.size = Pt(13)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_PRIMARY
    p_badge.space_after = Pt(14)

    p_main = tf1.add_paragraph()
    p_main.text = "AI-Based Smart Energy Consumption\nPrediction and Optimization System"
    p_main.font.size = Pt(36)
    p_main.font.bold = True
    p_main.font.color.rgb = COLOR_WHITE
    p_main.space_after = Pt(16)

    p_sub = tf1.add_paragraph()
    p_sub.text = "A Full-Stack Machine Learning Solution for Real-Time Electricity Forecasting, Peak Load Shifting, and Automated Energy Conservation"
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = COLOR_SECONDARY
    p_sub.space_after = Pt(24)

    p_meta = tf1.add_paragraph()
    p_meta.text = "Tech Stack: Python 3.11  •  Flask  •  Random Forest Regressor  •  Scikit-Learn  •  SQLite  •  Bootstrap 5  •  Chart.js"
    p_meta.font.size = Pt(13)
    p_meta.font.color.rgb = COLOR_GRAY

    # ----------------------------------------------------
    # SLIDE 2: Introduction & Motivation
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "1. Introduction & Background Motivation")

    add_card(slide2, 0.8, 1.8, 5.6, 5.0, "The Energy Crisis & Billing Gap", [
        "Rapid urban growth and proliferation of heavy appliances (ACs, geysers, EV charging) drive surging household electricity consumption.",
        "Traditional utility meters provide delayed feedback through monthly utility bills, preventing consumers from taking timely action.",
        "High electricity bills accumulate without clear visibility into which specific appliances or time windows caused the surge."
    ], COLOR_PRIMARY)

    add_card(slide2, 6.8, 1.8, 5.7, 5.0, "Why Machine Learning?", [
        "Electricity consumption is driven by complex weather patterns, human schedules, and appliance usage that simple rules cannot model accurately.",
        "Machine Learning allows predictive foresight: projecting energy demand before it occurs so consumers can adjust behaviors proactively.",
        "Empowers consumers to reduce utility bills by 15% to 30% while helping the national grid mitigate evening peak demand surges."
    ], COLOR_SECONDARY)

    # ----------------------------------------------------
    # SLIDE 3: Problem Statement
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "2. Problem Statement & Challenges")

    add_card(slide3, 0.8, 1.8, 3.6, 5.0, "Lack of Predictive Foresight", [
        "Consumers have no tool to estimate hourly/daily energy demand prior to running high-draw appliances.",
        "Unanticipated weather fluctuations (high heat/humidity) lead to unexpected billing shocks.",
        "Absence of proactive notifications before entering higher tariff tiers."
    ], COLOR_ACCENT)

    add_card(slide3, 4.8, 1.8, 3.6, 5.0, "Peak Grid Demand Strain", [
        "Heavy electrical grid strain occurs during evening peak hours (18:00 to 22:00).",
        "Coincident usage causes grid congestion, requiring costly fossil peaker plants and inflating tariffs.",
        "Lack of consumer incentives or visibility to shift loads to off-peak periods."
    ], RGBColor(244, 63, 94))

    add_card(slide3, 8.8, 1.8, 3.7, 5.0, "Vampire Power & Inaction", [
        "Standby 'vampire' power draw accounts for 8% to 10% of household electricity without providing utility.",
        "Standard energy-saving tips are generic, non-contextual, and lack financial quantification.",
        "No interactive tool to simulate tangible savings from energy conservation actions."
    ], COLOR_SECONDARY)

    # ----------------------------------------------------
    # SLIDE 4: Project Objectives
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "3. Project Objectives & Scope")

    add_card(slide4, 0.8, 1.8, 5.6, 5.0, "Machine Learning Objectives", [
        "Generate a physically-grounded residential energy dataset (4,320+ hourly records across 180 days).",
        "Train an explainable Random Forest Regressor to forecast energy consumption in kWh.",
        "Benchmark against a classical Linear Regression model.",
        "Evaluate model performance using MAE, MSE, RMSE, and R² Score metrics.",
        "Extract feature importance rankings to explain prediction decisions."
    ], COLOR_PRIMARY)

    add_card(slide4, 6.8, 1.8, 5.7, 5.0, "Software Engineering Objectives", [
        "Develop a modular, 3-tier MVC full-stack web application using Flask and SQLAlchemy.",
        "Implement secure user registration, session management, and password hashing (PBKDF2/SHA256).",
        "Build an interactive Energy Saving Simulator with real-time financial savings calculation.",
        "Design a responsive energy-themed UI with Bootstrap 5 and Chart.js visualizations.",
        "Implement progressive domestic utility tiered slab cost calculation."
    ], COLOR_SECONDARY)

    # ----------------------------------------------------
    # SLIDE 5: System Architecture & Workflow
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "4. System Architecture & High-Level Flow")

    add_card(slide5, 0.8, 1.8, 3.6, 5.0, "Presentation Layer", [
        "HTML5 / CSS3 Design System with Light/Dark Mode toggle.",
        "Bootstrap 5 responsive layouts.",
        "6 Interactive Chart.js widgets (Trend, Weekday, Peak, Scatter).",
        "Reactive Energy Saving Simulator sliders."
    ], COLOR_PRIMARY)

    add_card(slide5, 4.8, 1.8, 3.6, 5.0, "Controller & Logic Layer", [
        "Flask 3.1 WSGI Application with modular Route Blueprints.",
        "Data Preprocessing & Lag Feature Engineering Pipeline.",
        "Dynamic Rule-Based Optimization Engine.",
        "Tariff & Tiered Slab Billing Calculators."
    ], COLOR_SECONDARY)

    add_card(slide5, 8.8, 1.8, 3.7, 5.0, "ML & Storage Layer", [
        "Random Forest Regressor serialized via Joblib (energy_model.pkl).",
        "SQLite Relational Database (energy.db).",
        "SQLAlchemy ORM (User, Prediction, EnergyRecord, Alert models).",
        "Data Ingestion & Cleansing Engine."
    ], COLOR_DARK)

    # ----------------------------------------------------
    # SLIDE 6: Dataset & Feature Engineering
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "5. Dataset Characteristics & Feature Engineering")

    add_card(slide6, 0.8, 1.8, 5.6, 5.0, "Synthesized Dataset (4,320 Records)", [
        "Spans 180 continuous days with hourly granularity (4,320 records).",
        "Physics-grounded thermal cooling demand: non-linear AC surge when temperature > 26°C.",
        "Diurnal behavioral cycles: morning wake-up surge (07:00–09:00), workday lull, evening peak (18:00–22:00).",
        "Occupancy scaling: accounts for 1 to 5 occupants with weekday/weekend variance."
    ], COLOR_PRIMARY)

    add_card(slide6, 6.8, 1.8, 5.7, 5.0, "10 Engineered Predictive Features", [
        "Time Features: hour (0–23), day, month, day_of_week (0–6), is_weekend (0/1).",
        "Weather Features: ambient temperature (°C), relative humidity (%).",
        "Household Features: number of people, appliance usage intensity (Low, Medium, High).",
        "Lag & Momentum Features: previous hour consumption (kWh), 3-hour rolling average (kWh)."
    ], COLOR_SECONDARY)

    # ----------------------------------------------------
    # SLIDE 7: Machine Learning Algorithm Selection
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "6. Machine Learning Algorithm: Random Forest Regressor")

    add_card(slide7, 0.8, 1.8, 5.6, 5.0, "Why Random Forest?", [
        "Ensemble Learning: Combines predictions of 100 decision trees via bagging (bootstrap aggregation) to drastically reduce variance.",
        "Non-Linear Modeling: Accurately captures non-linear thermal thresholds (e.g. AC compressor kicks in above 28°C) without complex manual transformations.",
        "Resistant to Overfitting: Constrained with max_depth=14 and min_samples_split=4.",
        "Explainable AI: Produces transparent Gini-based feature importance scores."
    ], COLOR_PRIMARY)

    add_card(slide7, 6.8, 1.8, 5.7, 5.0, "Why Not Deep Learning / Linear Regression?", [
        "Linear Regression Limitation: Assumes linear relationships, failing to capture diurnal peaks and thermal thresholds.",
        "Deep Learning / LSTM Limitation: Requires massive training datasets, GPU hardware, long training times, and acts as a 'black-box' with poor explainability.",
        "Random Forest is lightweight, CPU-trainable in seconds, and ideal for tabular residential energy forecasting."
    ], COLOR_ACCENT)

    # ----------------------------------------------------
    # SLIDE 8: Model Performance & Benchmark Results
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "7. Model Evaluation & Benchmark Comparison")

    add_card(slide8, 0.8, 1.8, 5.6, 5.0, "Random Forest (Primary Model)", [
        "R² Score: 0.9802 (98.02% variance explained)",
        "Mean Absolute Error (MAE): 0.2266 kWh",
        "Mean Squared Error (MSE): 0.1023",
        "Root Mean Squared Error (RMSE): 0.3198 kWh",
        "High accuracy with low error on 20% holdout test data (864 test samples)."
    ], COLOR_PRIMARY)

    add_card(slide8, 6.8, 1.8, 5.7, 5.0, "Linear Regression (Baseline Benchmark)", [
        "R² Score: 0.8301 (83.01% variance explained)",
        "Mean Absolute Error (MAE): 0.7069 kWh",
        "Mean Squared Error (MSE): 0.8759",
        "Root Mean Squared Error (RMSE): 0.9359 kWh",
        "Performance Delta: Random Forest improves R² by +15.01% and reduces MAE error by 67.9%."
    ], COLOR_SECONDARY)

    # ----------------------------------------------------
    # SLIDE 9: Model Explainability (Feature Importance)
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "8. Model Explainability & Feature Importance")

    add_card(slide9, 0.8, 1.8, 5.6, 5.0, "Feature Importance Ranking", [
        "1. Previous Consumption (Lag 1h): 57.58%",
        "2. Appliance Usage Intensity: 18.45%",
        "3. Hour of Day (Diurnal Curve): 15.06%",
        "4. Ambient Temperature: 4.26%",
        "5. Month (Seasonality): 1.33%",
        "6. Rolling 3h Average: 1.12%",
        "7. Relative Humidity: 0.94%",
        "8. Occupant Count: 0.82%"
    ], COLOR_PRIMARY)

    add_card(slide9, 6.8, 1.8, 5.7, 5.0, "Viva Defense Insights", [
        "Why Lag is #1: Energy usage exhibits strong autoregressive momentum (appliances running now tend to stay on in the next hour).",
        "Appliance & Hour Impact: Peak evening activities (cooking, lights, entertainment) account for over 33% of predictive power.",
        "Thermal Sensitivity: Outdoor temperature dictates AC compressor activation, heavily influencing summer spikes."
    ], COLOR_SECONDARY)

    # ----------------------------------------------------
    # SLIDE 10: Interactive Prediction & Cost Calculator
    # ----------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "9. Interactive Prediction & Cost Calculation")

    add_card(slide10, 0.8, 1.8, 5.6, 5.0, "Real-Time Prediction Interface", [
        "User inputs: Date, Time, Temperature, Humidity, Occupants, Appliance Level, Previous kWh.",
        "Instant Output: Predicted energy in kWh, hourly cost, daily projection, and monthly bill.",
        "Usage Categorization: Automated classification into Low (<3.5 kWh), Normal (3.5–7.5 kWh), or High (>7.5 kWh) with colored badges.",
        "Automated Logging: Predictions saved to SQLite database with audit timestamps."
    ], COLOR_PRIMARY)

    add_card(slide10, 6.8, 1.8, 5.7, 5.0, "Progressive Domestic Tier Billing", [
        "Flat Tariff Calculation: Cost = kWh × Tariff Rate (e.g. 10 kWh @ ₹8/kWh = ₹80).",
        "Progressive Domestic Tier Slabs:",
        "  • Tier 1 (0 – 100 kWh): 60% of base tariff (Lifeline Block)",
        "  • Tier 2 (101 – 300 kWh): 100% of base tariff (Standard Block)",
        "  • Tier 3 (> 300 kWh): 135% of base tariff (Heavy Surcharge)",
        "Multi-currency support: ₹ (INR default), $, €, £."
    ], COLOR_SECONDARY)

    # ----------------------------------------------------
    # SLIDE 11: Dynamic Optimization & Recommendations
    # ----------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    add_header(slide11, "10. Dynamic Optimization & Recommendation Engine")

    add_card(slide11, 0.8, 1.8, 5.6, 5.0, "Context-Aware Rule Processing", [
        "Peak Load Shifting: Detects 18:00–22:00 window and advises deferring washing machines/pumps to morning off-peak hours (saves ~₹432/mo).",
        "Thermostat Optimization: Advises setting AC to 24°C–26°C when temperature > 28°C (every 1°C saves ~6% cooling power).",
        "Humidity Control: Recommends AC 'Dry Mode' when humidity > 70% (saves up to 25% compressor energy).",
        "Load Staggering: Prevents simultaneous high-wattage equipment spikes."
    ], COLOR_PRIMARY)

    add_card(slide11, 6.8, 1.8, 5.7, 5.0, "Quantified Savings Output", [
        "Each recommendation dynamically computes both monthly kWh saved and money saved in user currency (₹).",
        "Eliminating Standby Vampire Power: Advises smart power strips during late night (saves ~₹160/mo).",
        "LED Upgrades: Lighting optimization advice tailored to occupant counts.",
        "Dynamic generation ensures advice reflects real-time input parameters rather than static hardcoded text."
    ], COLOR_SECONDARY)

    # ----------------------------------------------------
    # SLIDE 12: Interactive Energy Saving Simulator
    # ----------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    add_header(slide12, "11. Interactive Energy Saving Simulator")

    add_card(slide12, 0.8, 1.8, 5.6, 5.0, "Interactive User Controls", [
        "Reduce AC / Cooling Usage Slider (0% to 50%)",
        "Optimize Lighting Usage Slider (0% to 50%)",
        "Reduce Major Appliances Slider (0% to 50%)",
        "Shift Usage Away from Peak Hours Slider (0% to 50%)",
        "Eliminate Standby Vampire Power Toggle (Checkbox)",
        "Configurable Baseline Monthly kWh and Tariff Rate."
    ], COLOR_PRIMARY)

    add_card(slide12, 6.8, 1.8, 5.7, 5.0, "Real-Time Calculated Outputs", [
        "Live AJAX recalculation without page reload.",
        "Estimated Monthly & Yearly Monetary Savings (₹).",
        "Total Monthly Energy Reduced (kWh).",
        "New Projected Monthly Electricity Bill (₹).",
        "Animated progress bar displaying Overall Reduction Percentage (%).",
        "Detailed breakdown across cooling, lighting, appliances, and standby power."
    ], COLOR_SECONDARY)

    # ----------------------------------------------------
    # SLIDE 13: Dataset Upload & Visual Analytics
    # ----------------------------------------------------
    slide13 = prs.slides.add_slide(blank_layout)
    add_header(slide13, "12. Dataset Upload & Visual Analytics")

    add_card(slide13, 0.8, 1.8, 5.6, 5.0, "Dataset Upload & Cleansing", [
        "CSV File Upload with validation for headers, datatypes, and file size (<16MB).",
        "Automated data cleansing: handles missing values via median and forward-fill imputation.",
        "Preview table displaying first 10 cleaned rows and statistical summaries.",
        "Bulk ingestion into SQLite database with automatic chart updates.",
        "Downloadable Sample CSV Template (/download-sample-csv)."
    ], COLOR_PRIMARY)

    add_card(slide13, 6.8, 1.8, 5.7, 5.0, "6 Interactive Chart.js Visualizations", [
        "1. Daily Timeline Trend: Line chart showing daily usage over time.",
        "2. Weekly Analysis: Monday through Sunday comparative bar chart.",
        "3. Monthly Aggregate: Month-by-month historical consumption.",
        "4. Peak Hours Profile: 24-hour diurnal profile highlighting peak 18–22 window.",
        "5. Appliance Doughnut: Percentage breakdown across household loads.",
        "6. Temperature vs Energy: Scatter plot illustrating thermal sensitivity."
    ], COLOR_SECONDARY)

    # ----------------------------------------------------
    # SLIDE 14: Testing, Security & Verification
    # ----------------------------------------------------
    slide14 = prs.slides.add_slide(blank_layout)
    add_header(slide14, "13. Testing, Security & Verification")

    add_card(slide14, 0.8, 1.8, 5.6, 5.0, "Automated Pytest Suite (7/7 Passed)", [
        "Unit testing of data cleaning and feature extraction (TC-01).",
        "ML prediction inference and cost calculation testing (TC-02).",
        "Progressive tiered billing mathematics verification (TC-03).",
        "Energy Saving Simulator calculation verification (TC-04).",
        "End-to-end user registration, login, and session guard testing (TC-05).",
        "CSV upload schema validation & empty file handling (TC-06)."
    ], COLOR_PRIMARY)

    add_card(slide14, 6.8, 1.8, 5.7, 5.0, "Security & Error Handling", [
        "Password Security: PBKDF2/SHA-256 password hashing with salt (no plaintext passwords).",
        "Session Protection: Unauthorized route access redirected to login.",
        "Input Sanitization: Temperature, humidity, and kWh validated against physical bounds.",
        "SQL Injection Prevention: SQLAlchemy ORM parameterized queries.",
        "Graceful Error Handling: Custom 404, 500, and 413 error handlers."
    ], COLOR_SECONDARY)

    # ----------------------------------------------------
    # SLIDE 15: Conclusion & Future Scope
    # ----------------------------------------------------
    slide15 = prs.slides.add_slide(blank_layout)
    add_header(slide15, "14. Conclusion & Future Enhancements")

    add_card(slide15, 0.8, 1.8, 5.6, 5.0, "Project Conclusion", [
        "Successfully developed an intelligent, full-stack AI energy prediction and optimization system.",
        "High Accuracy: Random Forest Regressor achieves 98.02% R² score with MAE = 0.2266 kWh.",
        "Tangible Impact: Empowers users to achieve 15% to 30% monthly electricity bill savings.",
        "Explainable & Accessible: Completely offline, requiring no paid APIs or GPU hardware."
    ], COLOR_PRIMARY)

    add_card(slide15, 6.8, 1.8, 5.7, 5.0, "Future Scope", [
        "IoT Smart Meter Telemetry: Integrate ESP32 / PZEM-004T CT current sensors for real-time live streaming.",
        "Automated Smart Plug Actuation: Automatically turn off non-essential loads during peak tariffs.",
        "Rooftop Solar PV Integration: Forecast solar generation and optimize battery storage charging.",
        "Mobile App / PWA: Progressive Web App with automated WhatsApp/SMS budget alerts."
    ], COLOR_SECONDARY)

    # Save presentation
    prs.save(output_filename)
    print(f"Presentation saved successfully to: {output_filename}")
    return output_filename

if __name__ == "__main__":
    create_presentation()
